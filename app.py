from flask import Flask, render_template, request
from google import genai
from dotenv import load_dotenv
import os
import json
import PyPDF2
import io
import docx
import openpyxl
from pptx import Presentation

load_dotenv()

app = Flask(__name__)

client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

def extract_text(file, filename):
    ext = filename.rsplit('.', 1)[-1].lower()

    if ext == 'pdf':
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text

    elif ext == 'docx':
        doc = docx.Document(io.BytesIO(file.read()))
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == 'txt':
        return file.read().decode('utf-8')

    elif ext == 'xlsx':
        wb = openpyxl.load_workbook(io.BytesIO(file.read()))
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell]) + "\n"
        return text

    elif ext == 'pptx':
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    else:
        return ""

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    mode = request.form["mode"]
    count = request.form.get("count", "5")
    if not count:
        count = "5"

    if 'upload_file' in request.files and request.files['upload_file'].filename != '':
        uploaded = request.files['upload_file']
        notes = extract_text(uploaded, uploaded.filename)
        if not notes.strip():
            return render_template("index.html", error="Could not extract text from this file!", count=count)
    else:
        notes = request.form["notes"]

    if not notes.strip():
        return render_template("index.html", error="Please paste some notes or upload a file!", count=count)

    if mode == "summary":
        prompt = f"Summarize these notes clearly and concisely:\n\n{notes}"
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )
        return render_template("index.html", summary=response.text, mode=mode, count=count)

    elif mode == "flashcards":
        prompt = f"""Create {count} flashcards from these notes.
Return ONLY a JSON array like this, nothing else:
[
  {{"question": "What is X?", "answer": "X is..."}},
  {{"question": "What is Y?", "answer": "Y is..."}}
]

Notes:
{notes}"""
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )
        clean = response.text.strip().replace("```json", "").replace("```", "")
        flashcards = json.loads(clean)
        return render_template("index.html", flashcards=flashcards, mode=mode, count=count)

    elif mode == "quiz":
        prompt = f"""Create {count} multiple choice questions from these notes.
Return ONLY a JSON array like this, nothing else:
[
  {{
    "question": "What is X?",
    "options": ["A", "B", "C", "D"],
    "answer": "A"
  }}
]

Notes:
{notes}"""
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )
        clean = response.text.strip().replace("```json", "").replace("```", "")
        quiz = json.loads(clean)
        return render_template("index.html", quiz=quiz, mode=mode, count=count)

if __name__ == "__main__":
    app.run(debug=True)